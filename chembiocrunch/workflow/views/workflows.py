from django.shortcuts import render, redirect, get_object_or_404
from django.views.generic import View,  DetailView, ListView, CreateView
from django.views.generic.detail import SingleObjectMixin
from django.conf import settings

from braces.views import LoginRequiredMixin

from django.db.models import get_model
from workflow.forms import CreateWorkflowForm, IC50UploadForm, DataMappingForm, DataMappingFormSetHelper, DataMappingFormSet, ResetButton, VisualisationForm, HeatmapForm
from django.core.urlresolvers import reverse
from crispy_forms.layout import Layout, Div, Submit, HTML, Button, Row, Field, Fieldset, Reset
from django.http import HttpResponseRedirect, HttpResponse

import matplotlib

import matplotlib.pyplot as plt

import json

import seaborn as sns
from workflow.models import GRAPH_MAPPINGS
from seaborn import plotting_context, set_context
import mpld3
import pandas as pd
from django.forms import Form
from pptx import Presentation
from pptx.util import Inches, Px

class WorkflowView(LoginRequiredMixin):

    model = get_model("workflow", "Workflow")

    def get_queryset(self):
        '''Make sure that all of the views cannot see the object unless they own it!'''
        return self.model.objects.get_user_records(self.request.user)

    def get_context_data(self, **kwargs):
        # Call the base implementation first to get a context
        context = super(WorkflowView, self).get_context_data(**kwargs)
        context['revisions'] = [["upload" ,"not-done"],
                                ["validate" ,"not-done"],
                                ["visualise", "not-done"],
                                ["customise" , "not-done"]]
        return context

class IC50WorkflowView(LoginRequiredMixin):

    model = get_model("workflow", "IC50Workflow")

    def get_queryset(self):
        '''Make sure that all of the views cannot see the object unless they own it!'''
        return self.model.objects.get_user_records(self.request.user)

    def get_context_data(self, **kwargs):
        # Call the base implementation first to get a context
        context = super(IC50WorkflowView, self).get_context_data(**kwargs)
        context['revisions'] = [["upload" ,"not-done"],
                                ["validate" ,"not-done"],
                                ["visualise", "not-done"],
                                ["customise" , "not-done"]]
        return context


class WorkflowListView(WorkflowView, ListView):
    '''Lists only the workflows that belong to that user'''
    template_name = "workflows/workflow_list.html"

    def get_context_data(self, **kwargs):
        # Call the base implementation first to get a context
        context = super(WorkflowListView, self).get_context_data(**kwargs)
        # Add in a QuerySet of all the books

        context["object_list"] = list(context['object_list'])

        return context


    def post(self, request, *args, **kwargs):
        pass




class WorkflowCreateView(WorkflowView, CreateView ):
    '''creates a single workflow'''
    fields = ['title', 'uploaded_file']
    template_name = "workflows/workflow_create.html"
    form_class = CreateWorkflowForm
    success_url = 'success'

    def get_success_url(self):
        return reverse('workflow_data_mappings_edit', kwargs={
                'pk': self.object.pk,
                })

    def form_valid(self, form):
        user = self.request.user
        form.instance.created_by = user
        form_valid = super(WorkflowCreateView, self).form_valid(form)
        
        if get_model("workflow", "workflowdatamappingrevision").objects.get_mapping_revisions_for_workflow(self.object).count() == 0:
            self.object.create_first_data_revision()

        return form_valid

    def get_context_data(self, **kwargs):
        # Call the base implementation first to get a context
        context = super(WorkflowCreateView, self).get_context_data(**kwargs)
        context['revisions'][0][1] = "in-progress"
        return context



class IC50WorkflowCreateView(IC50WorkflowView, CreateView ):
    '''creates a single workflow'''
    fields = ['title', 'uploaded_data_file','uploaded_config_file', 'uploaded_meta_file']
    template_name = "workflows/workflow_ic50_create.html"
    #form_class = LabcyteEchoIC50UploadForm
    form_class = IC50UploadForm
    success_url = 'success'

    model = get_model("workflow", "IC50Workflow")

    def get_success_url(self):
        return reverse('workflow_ic50_heatmap', kwargs={
                'pk': self.object.pk,
                })

    def form_valid(self, form):
        user = self.request.user
        form.instance.created_by = user
        form_valid = super(IC50WorkflowCreateView, self).form_valid(form)

        return form_valid

    def get_context_data(self, **kwargs):
        # Call the base implementation first to get a context
        context = super(IC50WorkflowCreateView, self).get_context_data(**kwargs)
        context['revisions'][0][1] = "in-progress"
        return context




class WorkflowDetailView(WorkflowView, DetailView, ):
    pass

class IC50WorkflowDetailView(IC50WorkflowView, DetailView, ):
    pass


class WorkflowOperationChooser(WorkflowDetailView):
    '''Allows the user to choose what to do with their data'''
    def get(self, request, *args, **kwargs):
        return render(request,
            "workflows/workflow_operation_chooser.html")

    def post(self, request, *args, **kwargs):
        pass


def vis_label_filter_prexix(vis_id):
    return "vis_label_filter_%d" % vis_id

def vis_option_prexix(vis_id):
    return "vis_options_%d" % vis_id







class WorkflowDataMappingEditView(WorkflowDetailView):
    '''Allows the user to edit the data mappings created automatically for the data they import'''
    template_name = "workflows/workflow_data_mapping_edit.html"
    def get_context_data(self, **kwargs):
        # Call the base implementation first to get a context
        context = super(WorkflowDataMappingEditView, self).get_context_data(**kwargs)
        if not "formset" in kwargs:
            context["formset"] = DataMappingFormSet(initial=self.object.get_data_mapping_formset_data(), prefix="data_mappings")
        context.update(kwargs)
        helper = DataMappingFormSetHelper()
        helper.template = 'table_inline.html'
        helper.add_input(Submit("submit","Save data mappings and visualise"))
        helper.add_input(ResetButton("reset", "Reset Form"))

        context["helper"] = helper
        visualisations = get_model("workflow","visualisation").objects.by_workflow(self.object)
        # vis_dicts = []
        # for vis in visualisations:
        #     initial_data = vis.get_initial_form_data()
        #     column_data = vis.get_column_form_data()
        #     vis_dict = vis.__dict__
        #     vis_dict["vis_option_form"] = VisualisationForm(initial=initial_data, prefix=vis_option_prexix(vis.id), column_data=column_data)
        #     vis_dicts.append(vis_dict)
        # context["visualisations"] = vis_dicts

        if visualisations.count() == 0:
            context['revisions'][1][1] = "in-progress"
            context['revisions'][2][1] = "not-done"
        else:
            context['revisions'][1][1] = "done"
            context['revisions'][2][1] = "in-progress"
        context['revisions'][0][1] = "done"
        return context

    def post(self, request, *args, **kwargs):
        '''This view will always add a new graph, graph updates are handled by ajax'''
        self.object = self.get_object()
        formset = DataMappingFormSet(request.POST, prefix="data_mappings")
        if formset.is_valid():
            workflow_revision = formset.process(self.object)

            return HttpResponseRedirect(reverse("visualisation_builder",kwargs={
                'pk': self.object.id,
                'workflow_revision_id' : workflow_revision.id,
                }))
        return self.render_to_response(self.get_context_data(formset=formset))


class WorkflowHeatmapView(IC50WorkflowDetailView):
    
    template_name = "workflows/workflow_ic50_heatmap.html"

    

    
    def get_context_data(self, **kwargs):
        # Call the base implementation first to get a context
        context = super(WorkflowHeatmapView, self).get_context_data(**kwargs)
        form_class = HeatmapForm
        #steps_json = json.loads(self.object.get_latest_workflow_revision().steps_json)
        workflow_revision = self.object.get_latest_workflow_revision()
        steps_json = json.loads(workflow_revision.steps_json)
        
        context["heatmap_form"] = HeatmapForm(uploaded_data=self.object.get_data(), steps_json=steps_json)
        context['revisions'][0][1] = "done"
        context['revisions'][1][1] = "in-progress"
        context.update(kwargs)
        return context

    def post(self, request, *args, **kwargs):
        '''This view will always add a new graph, graph updates are handled by ajax'''
        self.object = self.get_object()

        workflow_revision = self.object.get_latest_workflow_revision()        

        steps_json = json.loads(workflow_revision.steps_json)


        form = HeatmapForm(request.POST, uploaded_data=self.object.get_data(), steps_json=steps_json)
        if form.is_valid():
            #print form.cleaned_data
            steps_json = json.loads(workflow_revision.steps_json)
            #form.cleaned_data()
            #loop through the stored json and for each pair, set the value to that in the form data
            #heatmap_json = context["steps_json"]
            for key, value in steps_json.iteritems():
                steps_json[key] = form.cleaned_data[key]

            workflow_revision.steps_json = json.dumps(steps_json)
            workflow_revision.save()
            workflow_revision.create_ic50_data()
            visualisation_id = workflow_revision.visualisations.all()[0].id

            return HttpResponseRedirect(
                    reverse('ic50_update_view', kwargs={
                    'pk': self.object.pk,
                    'workflow_revision_id' : self.object.get_latest_workflow_revision().id,
                    "visualisation_id" : visualisation_id
                    })
                )

class VisualisationBuilderView(WorkflowDetailView):
    '''Creates a new visualisation for a given workflow revision'''
    workflow_revision_id = None
    template_name = "visualise/visualisation_builder.html"

    def get(self, request, *args, **kwargs):
        '''Borrowed from django base detail view'''
        self.workflow_revision_id = kwargs.pop("workflow_revision_id")
        self.object = self.get_object()
        context = self.get_context_data(object=self.object)
        return self.render_to_response(context)

    def get_context_data(self,  **kwargs):
        if "visualisation_form" in kwargs:
            context["visualisation_form"] = kwargs.pop("visualisation_form")
        context = super(VisualisationBuilderView, self).get_context_data(**kwargs)      
        
        if self.workflow_revision_id == 'latest':
            context["workflow_revision"] = self.object.get_latest_workflow_revision()
        else:
            context["workflow_revision"] = get_object_or_404(get_model("workflow", "WorkflowDataMappingRevision"), pk=self.workflow_revision_id)
        context["visualisation_list"] = get_model("workflow", "visualisation").objects.by_workflow(self.object)
        if not context.get("visualisation_form", None):
            column_data = context["workflow_revision"].get_column_form_data()
            context["visualisation_form"] = VisualisationForm(column_data=column_data)

        context['revisions'][1][1] = "done"
        context['revisions'][2][1] = "in-progress"
        context['revisions'][0][1] = "done"
        holder_object = get_model("workflow", "visualisation")(x_axis=context["workflow_revision"].x_axis, 
                                                                y_axis=context["workflow_revision"].y_axis,
                                                                visualisation_type="bar",
                                                                data_mapping_revision=context["workflow_revision"])
        html = holder_object.get_svg()
        context['visualisation'] = {"html" : html}
        return context

    def post(self, request, *args, **kwargs):
        '''This view will always add a new graph, graph updates are handled by ajax'''
        
        self.workflow_revision_id = kwargs.pop("workflow_revision_id")
        self.object = self.get_object()
        if request.POST.get("delete", False) == "delete":        
            return self.render_to_response(self.get_context_data())
        if self.workflow_revision_id == 'latest':
            workflow_revision = self.object.get_latest_workflow_revision()
        else:
            workflow_revision = get_object_or_404(get_model("workflow", "WorkflowDataMappingRevision"), pk=self.workflow_revision_id)
        column_data = workflow_revision.get_column_form_data()
        vis_create_form = VisualisationForm(request.POST,column_data=column_data )
        if vis_create_form.is_valid():
            new_object = vis_create_form.save(workflow_revision)
            return HttpResponseRedirect(reverse("visualisation_update_view",kwargs={
                'pk': self.object.id,
                'workflow_revision_id' : workflow_revision.id,
                'visualisation_id' : new_object.id,
                }))
        return self.render_to_response(self.get_context_data(visualisation_form=vis_create_form))


class VisualisationUpdateView(WorkflowDetailView):

    workflow_revision_id = None
    visualisation_id = None
    template_name = "visualise/visualisation_builder.html"
    
    def get_context_data(self,  **kwargs):
        visualisation_form = None
        if "visualisation_form" in kwargs:
            visualisation_form = kwargs.pop("visualisation_form")

        context = super(VisualisationUpdateView, self).get_context_data(**kwargs)
        context["visualisation_id"] = self.visualisation_id
        context["workflow_revision"] = get_object_or_404(get_model("workflow", "WorkflowDataMappingRevision"), pk=self.workflow_revision_id)
        context["visualisation_list"] = get_model("workflow", "visualisation").objects.by_workflow(self.object)
        context["visualisation"] = context["visualisation_list"].get(pk=self.visualisation_id)

        context['revisions'][1][1] = "done"
        context['revisions'][2][1] = "done"
        context['revisions'][0][1] = "done"
        context['revisions'][3][1] = "in-progress"

        if not visualisation_form:
            column_data = context["visualisation"].get_column_form_data()
            context["visualisation_form"] = VisualisationForm(column_data=column_data, prefix="")
        else:
            context["visualisation_form"] = visualisation_form
        return context


    def get(self, request, *args, **kwargs):
        '''Borrowed from django base detail view'''
        self.workflow_revision_id = kwargs.pop("workflow_revision_id")
        self.visualisation_id = kwargs.pop("visualisation_id")

        self.object = self.get_object()
        context = self.get_context_data(object=self.object)
        return self.render_to_response(context)


    def post(self, request, *args, **kwargs):

        self.workflow_revision_id = kwargs.pop("workflow_revision_id")
        self.visualisation_id = kwargs.pop("visualisation_id")

        self.object = self.get_object()
        workflow_revision = get_object_or_404(get_model("workflow", "WorkflowDataMappingRevision"), pk=self.workflow_revision_id)
        visualisation_list = get_model("workflow", "visualisation").objects.by_workflow_revision(workflow_revision)
        visualisation = visualisation_list.get(pk=self.visualisation_id)
        if request.POST.get("delete", False) == "delete":
            visualisation.delete()
            return HttpResponseRedirect(reverse("visualisation_builder",kwargs={
                'pk': self.object.id,
                'workflow_revision_id' : workflow_revision.id,
                })) 
        column_data = visualisation.get_column_form_data()
        vis_update_form = VisualisationForm(request.POST,column_data=column_data )
        if vis_update_form.is_valid():
            new_object = vis_update_form.save(workflow_revision, visualisation=visualisation)
            context = self.get_context_data(object=self.object, visualisation_form=vis_update_form)
            print request.POST
            if request.POST.get("new", False) == "new":
                return HttpResponseRedirect(reverse("visualisation_builder",kwargs={
                'pk': self.object.id,
                'workflow_revision_id' : workflow_revision.id,
                })) 
            return self.render_to_response(context)
        else:
            print vis_update_form.errors




# class VisualisationView(DetailView):

#     template_name = "visualise/vis_new.html"
#     model = get_model("workflow", "visualisation")
    
#     def get_context_data(self, **kwargs):
#         # Call the base implementation first to get a context
#         context = super(VisualisationView, self).get_context_data(**kwargs)
#         if "html" in kwargs:
#             context["html"] = kwargs["html"]
#         return context

#     def get(self, request, *args, **kwargs):
        
#         self.object = self.get_object()
#         x = self.object.x_axis
#         y = self.object.y_axis
#         vis = self.object.visualisation_type

#         df = self.object.data_mapping_revision.get_data()

#         with plotting_context( "poster" ):

#             g = sns.FacetGrid(df, size=10, aspect=2)
            
#             g.map(GRAPH_MAPPINGS[vis]["function"], x, y);
#             #g.map(GRAPH_MAPPINGS[self.object.graph_type]["function"], self.object.x_axis, self.object.y_axis);
#             matplotlib.pyplot.plot()
#             fc = matplotlib.backends.backend_agg.FigureCanvasAgg(plt.figure(1))
#             response = HttpResponse(content_type='image/png')
#             fc.print_png(response)
#             return response


class VisualisationView(DetailView,):
    model = get_model("workflow", "visualisation")
    format = None
    def get(self, request, *args, **kwargs):
        self.object = self.get_object()
        self.format = kwargs.pop("format")
        fig = self.object.get_fig_for_dataframe()
        self.fc = matplotlib.backends.backend_agg.FigureCanvasAgg(fig)
        if self.format=="png":
            return self.get_png()
        if self.format=="svg":
            return self.get_svg()       
        if self.format=="ppt":
            return self.get_ppt()   

    def get_png(self):
        response = HttpResponse(content_type='image/png')
        self.fc.print_png(response,  transparent=True)
        response['Content-Disposition'] = 'attachment; filename="filename.png"'
        return response



    def get_svg(self):
        response = HttpResponse(self.object.html, content_type='image/svg+xml')
        response['Content-Disposition'] = 'attachment; filename="filename.svg"'
        return response


    def get_ppt(self):
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        self.fc.print_png("/tmp/test.png",  transparent=True)
        left =  Inches(0)
        top = Inches(1.5)
        pic = slide.shapes.add_picture("/tmp/test.png", left, top)
        file_path = '/tmp/test.pptx'
        prs.save(file_path)
        fsock = open(file_path,"r")
        response = HttpResponse(fsock, content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        response['Content-Disposition'] = 'attachment; filename=TestPpt.pptx'
        #fc.print_png(response)
        return response







# class VisualisationExportView(WorkflowView, DetailView,):
#     model = get_model("workflow", "visualisation")
    
#     def get(self, request, *args, **kwargs):
#         #self.self.object = self.get_object()
#         #df = self.object.data_mapping_revision.get_data()
#         with plotting_context( "talk" ):

#             #g = sns.FacetGrid(df, size=10, aspect=2)
#             #g.map(GRAPH_MAPPINGS[self.object.graph_type]["function"], self.object.x_axis, self.object.y_axis);
#             #g.map(GRAPH_MAPPINGS["bar"]["function"], self.object.x_axis, self.object.y_axis)
#             #plt.plot()
#             #fc = FigureCanvas(plt.figure(1))

#             #from here, send this image to python-pptx
#             #img_path = 'monty-truth.png'








class Ic50UpdateView(IC50WorkflowDetailView):

    workflow_revision_id = None
    visualisation_id = None
    template_name = "visualise/visualisation_builder.html"
    
    def get_context_data(self,  **kwargs):
        # visualisation_form = None
        # if "visualisation_form" in kwargs:
        #     visualisation_form = kwargs.pop("visualisation_form")

        context = super(Ic50UpdateView, self).get_context_data(**kwargs)
        context["workflow_revision"] = get_object_or_404(get_model("workflow", "IC50WorkflowRevision"), pk=self.workflow_revision_id)
        visualisation_id = self.visualisation_id
        if visualisation_id == 0:
            #This will return the first of the IC50 curves but also generate the others
            context["workflow_revision"].create_ic50_curves()
            visualisation_id = context["workflow_revision"].visualisations.all()[0].id
            


        context["visualisation_id"] = visualisation_id
        context["visualisation_list"] = context["workflow_revision"].visualisations.all()
        context["visualisation"] = context["visualisation_list"].get(pk=visualisation_id)

        context['revisions'][1][1] = "done"
        context['revisions'][2][1] = "done"
        context['revisions'][0][1] = "done"
        context['revisions'][3][1] = "in-progress"
        context["isic50"] = True

        # if not visualisation_form:
        #     column_data = context["visualisation"].get_column_form_data()
        #     context["visualisation_form"] = VisualisationForm(column_data=column_data, prefix="")
        # else:
        #     context["visualisation_form"] = visualisation_form
        return context


    def get(self, request, *args, **kwargs):
        '''Borrowed from django base detail view'''
        self.workflow_revision_id = kwargs.pop("workflow_revision_id")
        self.visualisation_id = kwargs.pop("visualisation_id")

        self.object = self.get_object()
        context = self.get_context_data(object=self.object)
        return self.render_to_response(context)


    def post(self, request, *args, **kwargs):

        self.workflow_revision_id = kwargs.pop("workflow_revision_id")
        self.visualisation_id = kwargs.pop("visualisation_id")

        self.object = self.get_object()
        workflow_revision = get_object_or_404(get_model("workflow", "IC50Visualisation"), pk=self.workflow_revision_id)
        visualisation_list = get_model("workflow", "IC50Visualisation").objects.by_workflow_revision(workflow_revision)
        visualisation = visualisation_list.get(pk=self.visualisation_id)
        if request.POST.get("delete", False) == "delete":
            visualisation.delete()
            return HttpResponseRedirect(reverse("visualisation_builder",kwargs={
                'pk': self.object.id,
                'workflow_revision_id' : workflow_revision.id,
                })) 
        column_data = visualisation.get_column_form_data()
        vis_update_form = VisualisationForm(request.POST,column_data=column_data )
        if vis_update_form.is_valid():
            new_object = vis_update_form.save(workflow_revision, visualisation=visualisation)
            context = self.get_context_data(object=self.object, visualisation_form=vis_update_form)
            print request.POST
            if request.POST.get("new", False) == "new":
                return HttpResponseRedirect(reverse("visualisation_builder",kwargs={
                'pk': self.object.id,
                'workflow_revision_id' : workflow_revision.id,
                })) 
            return self.render_to_response(context)
        else:
            print vis_update_form.errors

